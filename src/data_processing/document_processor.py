from typing import Dict, List, Any, Optional, Union, Tuple
import logging
from pathlib import Path
import re
import json
from PIL import Image
import pytesseract
import docx
import fitz  # PyMuPDF
from pptx import Presentation
import openpyxl
import markdown
import mammoth
from pdf2image import convert_from_path
import cv2
import numpy as np
import io
import uuid
from datetime import datetime
from .image_annotator import ImageAnnotator
import bisect

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Processes various document formats and extracts text and images."""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize document processor with configuration.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        self.doc_config = config["agent"]["document_processing"]
        
        # Initialize ImageAnnotator
        self.image_annotator = ImageAnnotator(config)
        
        # Configure OCR and image captioning
        self._configure_ocr()
        if self.doc_config.get("use_image_captioning", False):
            self._initialize_image_captioner()
        
        # Initialize output directories
        self.base_dir = Path(config["paths"]["base_dir"])
        self.image_output_dir = self.base_dir / "data" / "processed" / "extracted_images"
        self.image_output_dir.mkdir(parents=True, exist_ok=True)
        
        # Create subdirectories for different sources
        (self.image_output_dir / "pdf").mkdir(exist_ok=True)
        (self.image_output_dir / "docx").mkdir(exist_ok=True)
        (self.image_output_dir / "pptx").mkdir(exist_ok=True)
        
        # Initialize image metadata tracking
        self.image_metadata = {}

    def _configure_ocr(self) -> None:
        """Configure OCR settings."""
        ocr_config = self.doc_config.get("ocr_config", {})
        self.ocr_enabled = self.doc_config.get("ocr_enabled", True)
        if self.ocr_enabled:
            self.ocr_config = {
                "lang": ocr_config.get("language", "eng"),
                "config": f"--psm {ocr_config.get('psm', 3)} --oem {ocr_config.get('oem', 3)}"
            }
    
    def _initialize_image_captioner(self) -> None:
        """Initialize image captioning model."""
        try:
            from transformers import pipeline
            model_name = self.doc_config.get("image_caption_model", "Salesforce/blip2-flan-t5-xl")
            self.image_captioner = pipeline("image-to-text", model=model_name)
        except Exception as e:
            logger.error(f"Failed to initialize image captioner: {e}")
            self.image_captioner = None
    
    def _save_image(self, image: Union[Image.Image, bytes], source_path: Path, location_info: Dict[str, Any] = None) -> Dict[str, Any]:
        """Save extracted image and return its metadata."""
        try:
            # Convert bytes to PIL Image if needed
            if isinstance(image, bytes):
                image = Image.open(io.BytesIO(image))
            
            # Generate unique image ID and filename
            image_id = str(uuid.uuid4())
            source_type = source_path.suffix.lower()[1:]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{source_path.stem}_{timestamp}_{image_id}.png"
            
            # Determine output directory based on source type
            output_dir = self.image_output_dir / source_type
            output_path = output_dir / filename
            
            # Save image
            image.save(output_path, "PNG")
            
            # Create metadata
            metadata = {
                "id": image_id,
                "source_document": str(source_path),
                "source_type": source_type,
                "filename": filename,
                "path": str(output_path),
                "size": image.size,
                "format": image.format,
                "extracted_timestamp": timestamp
            }
            
            # Add location info if available
            if location_info:
                metadata["location"] = location_info
            
            # Store metadata
            self.image_metadata[image_id] = metadata
            
            # Save metadata to file
            self._save_metadata()
            
            return metadata
            
        except Exception as e:
            logger.error(f"Error saving image from {source_path}: {e}")
            return {}

    def _save_metadata(self) -> None:
        """Save image metadata to JSON file."""
        try:
            metadata_file = self.image_output_dir / "image_metadata.json"
            with open(metadata_file, 'w') as f:
                json.dump(self.image_metadata, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving image metadata: {e}")

    def _process_image_content(self, image: Image.Image) -> str:
        """Process image using OCR and/or captioning."""
        extracted_text = []
        
        # Perform OCR if enabled
        if self.ocr_enabled:
            try:
                ocr_text = pytesseract.image_to_string(
                    image,
                    lang=self.ocr_config["lang"],
                    config=self.ocr_config["config"]
                )
                if ocr_text.strip():
                    extracted_text.append(f"[OCR Text: {ocr_text.strip()}]")
            except Exception as e:
                logger.warning(f"OCR failed: {e}")
        
        # Generate image caption if enabled
        if self.image_captioner:
            try:
                captions = self.image_captioner(image)
                if captions and isinstance(captions, list) and captions[0].get('generated_text'):
                    caption = captions[0]['generated_text']
                    extracted_text.append(f"[Image Caption: {caption}]")
            except Exception as e:
                logger.warning(f"Image captioning failed: {e}")
        
        return "\n".join(extracted_text) if extracted_text else ""

    def _extract_and_process_images(self, doc_images: List[Dict[str, Any]], source_type: str) -> str:
        """Extract and process images from document."""
        image_content = []
        
        for img_data in doc_images:
            try:
                if isinstance(img_data.get("image"), Image.Image):
                    image = img_data["image"]
                elif isinstance(img_data.get("image"), bytes):
                    image = Image.open(io.BytesIO(img_data["image"]))
                else:
                    continue
                    
                # Process image content using ImageAnnotator
                result = self.image_annotator.process_image(image)
                
                # Add location context if available
                location_info = f" at {img_data.get('location', 'unknown location')}" if 'location' in img_data else ""
                
                # Format image content with all extracted information
                content_parts = []
                
                if result["ocr_text"]:
                    content_parts.append(f"[OCR Text{location_info}: {result['ocr_text']}]")
                
                if result["caption"]:
                    content_parts.append(f"[Image Caption{location_info}: {result['caption']}]")
                
                # Add rich context information
                context = result.get("context", {})
                
                if context.get("objects"):
                    obj_text = ", ".join(
                        f"{obj['label']} ({obj['confidence']:.2f})"
                        for obj in context["objects"]
                    )
                    content_parts.append(f"[Detected Objects{location_info}: {obj_text}]")
                
                if context.get("scene_description"):
                    content_parts.append(
                        f"[Scene Description{location_info}: {context['scene_description']}]"
                    )
                
                if context.get("visual_qa"):
                    vqa_text = " | ".join(
                        f"{qa['question']} -> {qa['answer']}"
                        for qa in context["visual_qa"]
                        if qa["answer"].lower() not in ["unknown", "none", "not sure"]
                    )
                    if vqa_text:
                        content_parts.append(f"[Visual Analysis{location_info}: {vqa_text}]")
                
                if content_parts:
                    image_content.append("\n".join(content_parts))
                    
            except Exception as e:
                logger.warning(f"Failed to process image: {e}")
        
        return "\n\n".join(image_content)

    def _merge_text_and_image_content(self, text_content: str, image_content: str, relationships: List[Dict[str, Any]]) -> str:
        """Merge text and image content based on relationships and context."""
        if not image_content:
            return text_content
            
        if not relationships:
            # If no relationships defined, structure content in a meaningful way
            sections = []
            if text_content:
                sections.append("=== Document Text ===\n" + text_content)
            if image_content:
                sections.append("=== Visual Content ===\n" + image_content)
            return "\n\n".join(sections)
            
        # Convert text content to list of paragraphs
        paragraphs = text_content.split("\n\n")
        image_sections = image_content.split("\n\n")
        
        # Create a mapping of image content to their targets
        image_map = {}
        for rel, img_section in zip(relationships, image_sections):
            if rel["type"] in {"embedded_content", "page_content"}:
                if "context" in rel:
                    # Enhanced content with context mapping
                    image_map[rel["target"]] = {
                        "content": img_section,
                        "context": rel["context"],
                        "position": rel.get("position", "after")
                    }
                else:
                    image_map[rel["target"]] = {
                        "content": img_section,
                        "context": None,
                        "position": "after"
                    }
        
        # Insert image content at appropriate locations
        result = []
        for i, paragraph in enumerate(paragraphs):
            result.append(paragraph)
            
            # Check if any images should be inserted here
            for rel in relationships:
                if rel["source"] == f"paragraph_{i}" or rel["source"] == f"page_{i}":
                    if img_data := image_map.get(rel["target"]):
                        # Add contextual separator
                        if img_data["context"]:
                            result.append(f"\n[Visual Context: {img_data['context']}]")
                        result.append(img_data["content"])
                        
                        # Add visual-textual relationship marker if available
                        if img_data.get("context") and "Detected Objects" in img_data["content"]:
                            result.append(
                                "\n[This visual content provides additional context for the surrounding text]"
                            )
        
        return "\n\n".join(result)

    def _process_document_images(self, doc_images: List[Dict[str, Any]], source_type: str) -> Tuple[str, List[Dict[str, Any]]]:
        """Process images from document and generate relationships."""
        if not doc_images:
            return "", []
            
        image_results = []
        relationships = []
        
        for img_idx, img_data in enumerate(doc_images):
            try:
                # Process image using ImageAnnotator
                if isinstance(img_data.get("image"), Image.Image):
                    image = img_data["image"]
                elif isinstance(img_data.get("image"), bytes):
                    image = Image.open(io.BytesIO(img_data["image"]))
                else:
                    continue
                    
                # Extract comprehensive image content
                result = self.image_annotator.process_image(image)
                
                source = img_data.get("paragraph_index", img_idx)
                # Create relationship with enhanced context
                relationship = {
                    "type": "embedded_content",
                    "source": f"paragraph_{source}",
                    "target": f"img_{img_idx}",
                    "context": self._generate_image_context(result)
                }
                
                if "location" in img_data:
                    relationship["position"] = "at_location"
                    relationship["location"] = img_data["location"]
                
                relationships.append(relationship)
                image_results.append(result)
                
            except Exception as e:
                logger.warning(f"Failed to process image {img_idx}: {e}")
        
        # Combine all image results into a single text representation
        combined_text = self.image_annotator.combine_results(image_results)
        
        return combined_text, relationships
    
    def _generate_image_context(self, result: Dict[str, Any]) -> str:
        """Generate a concise context description from image analysis results."""
        context_parts = []
        
        # Add high-confidence detected objects
        if objects := result.get("context", {}).get("objects"):
            high_conf_objects = [
                obj["label"] for obj in objects 
                if obj["confidence"] > 0.7
            ]
            if high_conf_objects:
                context_parts.append(f"Contains: {', '.join(high_conf_objects)}")
        
        # Add scene description if available
        if scene_desc := result.get("context", {}).get("scene_description"):
            context_parts.append(f"Scene: {scene_desc}")
        
        # Add relevant visual QA insights
        if vqa_results := result.get("context", {}).get("visual_qa"):
            relevant_answers = [
                qa["answer"] for qa in vqa_results
                if qa["answer"].lower() not in ["unknown", "none", "not sure"]
                and qa["question"] in [
                    "What is happening in this image?",
                    "What is the main subject of this image?"
                ]
            ]
            if relevant_answers:
                context_parts.append(f"Content: {' | '.join(relevant_answers)}")
        
        return " | ".join(context_parts) if context_parts else "Visual content"

    def process_document(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a document and extract its content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing extracted text, images, and relationships
        """
        suffix = file_path.suffix.lower()
        
        try:
            if suffix in {".txt", ".md"}:
                doc_data = self._process_text_file(file_path)
            elif suffix in {".doc", ".docx"}:
                doc_data = self._process_word_doc(file_path)
            elif suffix == ".pdf":
                doc_data = self._process_pdf(file_path)
            elif suffix in {".ppt", ".pptx"}:
                doc_data = self._process_powerpoint(file_path)
            elif suffix in {".xls", ".xlsx"}:
                doc_data = self._process_excel(file_path)
            else:
                raise ValueError(f"Unsupported file format: {suffix}")

            # Process images and merge content
            image_content, relationships = self._process_document_images(doc_data["images"], suffix[1:].upper())
            unified_content = self._merge_text_and_image_content(
                doc_data["text"],
                image_content,
                relationships
            )
            
            return {
                "text": unified_content,  # Contains both text and processed image content
                "images": doc_data["images"],  # Keep original images for reference
                "relationships": relationships
            }
                
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            raise
    
    def _process_text_file(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text or markdown files."""
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
            
        # Convert markdown to text if needed
        if file_path.suffix.lower() == ".md":
            content = markdown.markdown(content)
            # Strip HTML tags
            content = re.sub(r"<[^>]+>", "", content)
        
        # Apply text settings
        content = self._clean_text(content)
        
        return {
            "text": content,
            "images": [],
            "relationships": []
        }
    
    def _process_word_doc(self, file_path: Path) -> Dict[str, Any]:
        """Process Word documents."""
        if file_path.suffix.lower() == ".docx":
            doc = docx.Document(file_path)
            text_content = []
            images = []
            relationships = []
            
            for i, paragraph in enumerate(doc.paragraphs):
                text_content.append(paragraph.text)
                
                # Extract inline images
                for run in paragraph.runs:
                    if run._element.findall(".//pic:pic", {"pic": "http://schemas.openxmlformats.org/drawingml/2006/picture"}):
                        try:
                            image_data = run._element.findall(".//a:blip", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})[0]
                            image_id = image_data.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                            if image_id:
                                image_path = f"word/media/{image_id}"
                                image_bytes = doc.part.related_parts[image_id].blob
                                
                                # Save image and get metadata
                                location_info = {
                                    "paragraph_index": i
                                }
                                
                                image_meta = self._save_image(image_bytes, file_path, location_info)
                                
                                if image_meta:
                                    images.append({
                                        "id": image_meta["id"],
                                        "image": Image.open(io.BytesIO(image_bytes)),
                                        "metadata": image_meta
                                    })
                                    
                                    relationships.append({
                                        "type": "embedded_content",
                                        "source": f"paragraph_{i}",
                                        "target": image_meta["id"]
                                    })
                                    
                        except Exception as e:
                            logger.warning(f"Failed to extract image from paragraph {i}: {e}")
            
            return {
                "text": "\n".join(text_content),
                "images": images,
                "relationships": relationships
            }
        else:
            # For .doc files, use mammoth for conversion
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                text = result.value
                # Strip HTML tags
                text = re.sub(r"<[^>]+>", "", text)
                return {
                    "text": self._clean_text(text),
                    "images": [],  # Old .doc format doesn't support reliable image extraction
                    "relationships": []
                }
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF documents with enhanced extraction using PyMuPDF."""
        text_content = []
        images = []
        relationships = []
        
        try:
            # Open PDF with PyMuPDF
            with fitz.open(file_path) as pdf:
                for page_num, page in enumerate(pdf):
                    # Get page rotation and handle different orientations
                    rotation = page.rotation
                    if rotation != 0:
                        page.set_rotation(0)
                    
                    # Extract text with detailed layout analysis
                    layout = page.get_text("dict", sort=True, flags=11)  # Use all layout analysis flags
                    page_text = []
                    
                    # Process blocks with layout awareness
                    for block in layout["blocks"]:
                        if block["type"] == 0:  # Text block
                            block_text = []
                            for line in block["lines"]:
                                # Get text with font info
                                spans = []
                                for span in line["spans"]:
                                    spans.append({
                                        "text": span["text"],
                                        "font": span["font"],
                                        "size": span["size"],
                                        "flags": span["flags"]  # Bold, italic etc
                                    })
                                
                                # Preserve formatting for structured content
                                line_text = ""
                                current_format = None
                                for span in spans:
                                    if self._is_heading(span):
                                        line_text += f"\n## {span['text']} ##\n"
                                    elif self._is_list_item(span):
                                        line_text += f"\n• {span['text']}"
                                    else:
                                        line_text += span["text"]
                                
                                if line_text.strip():
                                    block_text.append(line_text)
                            
                            # Handle tables
                            if self._is_table_block(block):
                                table_text = self._extract_table_content(block)
                                if table_text:
                                    block_text.append(table_text)
                            
                            if block_text:
                                page_text.append(" ".join(block_text))
                    
                    if page_text:
                        text_content.append("\n".join(page_text))
                    
                    # Extract images if enabled
                    if self.doc_config["pdf_settings"]["extract_images"]:
                        self._extract_pdf_images(page, page_num, pdf, images, relationships)
                        
                    # Fallback to OCR if needed
                    if not page_text and self.doc_config["pdf_settings"]["use_ocr_fallback"]:
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text = self._perform_ocr(img)
                        if ocr_text:
                            text_content.append(ocr_text)
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise
        
        return {
            "text": "\n\n".join(text_content),
            "images": images,
            "relationships": relationships
        }

    def _extract_pdf_images(self, page: fitz.Page, page_num: int, pdf: fitz.Document, 
                            images: List[Dict[str, Any]], relationships: List[Dict[str, Any]]) -> None:
        """Extract and filter images from PDF page."""
        try:
            # Get list of images on the page
            image_list = page.get_images(full=True)  # Get full image info
            
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]  # Image reference
                    base_image = pdf.extract_image(xref)
                    
                    if base_image:
                        # Convert to PIL Image for analysis
                        image_data = base_image["image"]
                        image = Image.open(io.BytesIO(image_data))
                        
                        # Skip if image is too small (likely an icon)
                        if image.size[0] < 100 or image.size[1] < 100:
                            continue
                            
                        # Calculate image complexity using edge detection
                        img_array = np.array(image.convert('L'))
                        edges = cv2.Canny(img_array, 100, 200)
                        complexity = np.sum(edges > 0) / (img_array.shape[0] * img_array.shape[1])
                        
                        # Skip if image is too simple (likely a logo or icon)
                        if complexity < 0.01:
                            continue
                            
                        # Skip if image appears to be a single color or gradient
                        colors = image.getcolors(maxcolors=256)
                        if colors is not None and len(colors) < 5:
                            continue
                        
                        # Get image location and size on page
                        bbox = page.get_image_bbox(xref)
                        if bbox:
                            location_info = {
                                "page": page_num,
                                "bbox": list(bbox),  # Convert rect to list
                                "width": bbox.width,
                                "height": bbox.height,
                                "area_ratio": (bbox.width * bbox.height) / (page.rect.width * page.rect.height)
                            }
                            
                            # Skip if image takes up too little page space
                            if location_info["area_ratio"] < 0.01:
                                continue
                            
                            # Save image and get metadata
                            image_meta = self._save_image(image, pdf.name, location_info)
                            
                            if image_meta:
                                images.append({
                                    "id": image_meta["id"],
                                    "image": image,
                                    "metadata": image_meta
                                })
                                
                                relationships.append({
                                    "type": "page_content",
                                    "source": f"page_{page_num}",
                                    "target": image_meta["id"],
                                    "position": "at_location",
                                    "location": location_info
                                })
                                
                except Exception as e:
                    logger.warning(f"Failed to extract image {img_index} from page {page_num}: {e}")
                    
        except Exception as e:
            logger.error(f"Error extracting images from page {page_num}: {e}")

    def _enhance_image_quality(self, image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results."""
        try:
            # Convert to numpy array
            img_array = np.array(image)
            
            # Convert to grayscale if needed
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
                
            # Apply adaptive thresholding
            binary = cv2.adaptiveThreshold(
                gray, 255,
                cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                cv2.THRESH_BINARY, 11, 2
            )
            
            # Denoise using median blur
            denoised = cv2.medianBlur(binary, 3)
            
            # Deskew if needed
            coords = np.column_stack(np.where(denoised > 0))
            angle = cv2.minAreaRect(coords)[-1]
            if angle < -45:
                angle = 90 + angle
            if abs(angle) > 0.5:
                (h, w) = denoised.shape[:2]
                center = (w // 2, h // 2)
                M = cv2.getRotationMatrix2D(center, angle, 1.0)
                rotated = cv2.warpAffine(
                    denoised, M, (w, h),
                    flags=cv2.INTER_CUBIC,
                    borderMode=cv2.BORDER_REPLICATE
                )
            else:
                rotated = denoised
                
            return Image.fromarray(rotated)
        except Exception as e:
            logger.error(f"Error enhancing image quality: {e}")
            return image

    def _is_heading(self, span: Dict[str, Any]) -> bool:
        """Detect if text span is a heading based on font properties."""
        return (span["size"] > 12 and span["flags"] & 2**3) or span["size"] > 14  # Check bold flag or large size

    def _is_list_item(self, span: Dict[str, Any]) -> bool:
        """Detect if text span is a list item."""
        text = span["text"].strip()
        return text.startswith(("•", "-", "*", "○", "▪", "◦")) or re.match(r"^\d+\.", text)

    def _is_table_block(self, block: Dict[str, Any]) -> bool:
        """Detect if block contains tabular data based on layout analysis."""
        if not block.get("lines"):
            return False
        
        # Check for consistent column alignment
        x_positions = set()
        for line in block["lines"]:
            for span in line["spans"]:
                x_positions.add(round(span["origin"][0], 1))
        
        return len(x_positions) > 2 and all(
            abs(x1 - x2) > 5 for x1, x2 in zip(sorted(x_positions)[:-1], sorted(x_positions)[1:])
        )

    def _extract_table_content(self, block: Dict[str, Any]) -> str:
        """Extract and format tabular content."""
        if not block.get("lines"):
            return ""
        
        # Collect column positions
        x_positions = set()
        for line in block["lines"]:
            for span in line["spans"]:
                x_positions.add(round(span["origin"][0], 1))
        x_positions = sorted(list(x_positions))
        
        # Build table rows
        table_rows = []
        for line in block["lines"]:
            row = [""] * len(x_positions)
            for span in line["spans"]:
                col_idx = bisect.bisect_left(x_positions, round(span["origin"][0], 1))
                if 0 <= col_idx < len(row):
                    row[col_idx] += span["text"].strip()
            if any(cell.strip() for cell in row):
                table_rows.append(row)
        
        # Format as markdown table
        if table_rows:
            header = "| " + " | ".join(["" for _ in range(len(x_positions))]) + " |"
            separator = "|" + "|".join(["---" for _ in range(len(x_positions))]) + "|"
            rows = [
                "| " + " | ".join(row) + " |"
                for row in table_rows
            ]
            return "\n".join([header, separator] + rows)
    
    def _process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint presentations."""
        prs = Presentation(file_path)
        text_content = []
        images = []
        relationships = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
                
                # Extract images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image_bytes = shape.image.blob
                        
                        # Save image and get metadata
                        location_info = {
                            "slide": slide_num
                        }
                        
                        image_meta = self._save_image(image_bytes, file_path, location_info)
                        
                        if image_meta:
                            images.append({
                                "id": image_meta["id"],
                                "image": Image.open(io.BytesIO(image_bytes)),
                                "metadata": image_meta
                            })
                            
                            relationships.append({
                                "type": "slide_content",
                                "source": f"slide_{slide_num}",
                                "target": image_meta["id"]
                            })
                            
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_num}: {e}")
            
            text_content.append("\n".join(slide_text))
        
        return {
            "text": "\n".join(text_content),
            "images": images,
            "relationships": relationships
        }
    
    def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel spreadsheets."""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_content = []
        
        for sheet in wb:
            sheet_content = []
            for row in sheet.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                if any(row_values):  # Skip empty rows
                    if self.doc_config["text_settings"]["table_format"] == "markdown":
                        sheet_content.append("| " + " | ".join(row_values) + " |")
                    else:
                        sheet_content.append("\t".join(row_values))
            
            if sheet_content:
                text_content.append(f"Sheet: {sheet.title}")
                if self.doc_config["text_settings"]["table_format"] == "markdown":
                    # Add markdown table header separator
                    sheet_content.insert(1, "| " + " | ".join(["---"] * len(sheet_content[0].split("|")[1:-1])) + " |")
                text_content.extend(sheet_content)
                text_content.append("")  # Add spacing between sheets
        
        return {
            "text": "\n".join(text_content),
            "images": [],
            "relationships": []
        }
    
    def _perform_ocr(self, image: Union[Image.Image, np.ndarray]) -> str:
        """Perform OCR on an image."""
        try:
            # Convert numpy array to PIL Image if needed
            if isinstance(image, np.ndarray):
                image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
            
            # Perform OCR
            text = pytesseract.image_to_string(
                image,
                lang=self.ocr_config["lang"],
                config=self.ocr_config["config"]
            )
            
            return self._clean_text(text)
            
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text based on settings."""
        if not text:
            return ""
            
        settings = self.doc_config["text_settings"]
        
        if settings["clean_whitespace"]:
            # Normalize whitespace
            text = re.sub(r"\s+", " ", text)
            text = re.sub(r"\n\s*\n", "\n\n", text)
            text = text.strip()
        
        if settings["remove_urls"]:
            # Remove URLs
            text = re.sub(r"http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+", "", text)
        
        # Enforce length limits
        if len(text) < settings["min_chars"]:
            return ""
        if len(text) > settings["max_chars"]:
            text = text[:settings["max_chars"]]
        
        return text